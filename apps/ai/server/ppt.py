import os
import subprocess
from langchain.docstore.document import Document
from pdfplumber import open as pdf_open
from pptx import Presentation
from pdf import PDF

class PPT():
    def process_file(cls, ppt_path):
        pdf_path = os.path.join(os.getcwd(), "multimodal/ppt_references", os.path.basename(ppt_path).replace('.pptx', '.pdf').replace('.ppt', '.pdf'))
        cls.convert_ppt_to_pdf(ppt_path)
        images_data = cls.convert_pdf_to_images(pdf_path)
        slide_texts = cls.extract_text_and_notes_from_ppt(ppt_path)
        processed_data = []

        for (image_path, page_num), (slide_text, notes) in zip(images_data, slide_texts):
            if notes:
                notes = "\n\nThe speaker notes for this slide are: " + notes

            # get image description with NeVA/DePlot
            image_description = " "
            if PDF.is_graph(image_path):
                image_description = PDF.process_graph(image_path)
            caption = slide_text + image_description + notes
            image_metadata = {
                    "x1":0,
                    "y1":0,
                    "x2":0,
                    "x3":0,
                    "source": f"{os.path.basename(ppt_path)}",
                    "image": image_path,
                    "caption": caption,
                    "type": "image",
                    "page_num": page_num
            }
            processed_data.append(Document(page_content = "This is a slide with the text: " + slide_text + image_description, metadata = image_metadata))

        return processed_data
    
    def convert_ppt_to_pdf(ppt_path):
        base_name = os.path.basename(ppt_path)
        ppt_name_without_ext = os.path.splitext(base_name)[0].replace(" ", "_")

        # Use the existing directory '../../ppt_references/'
        new_dir_path = os.path.abspath("multimodal/ppt_references")

        # Set the new PDF path in the existing directory
        pdf_path = os.path.join(new_dir_path, f"{ppt_name_without_ext}.pdf")

        # LibreOffice command to convert PPT to PDF
        command = [
            'libreoffice', 
            '--headless', 
            '--convert-to', 
            'pdf', 
            '--outdir', 
            new_dir_path, 
            ppt_path
        ]
        subprocess.run(command, check=True)

        return pdf_path
    
    def convert_pdf_to_images(pdf_path):
        doc = pdf_open(pdf_path)

        # Extract the base name of the PDF file and replace spaces with underscores
        base_name = os.path.basename(pdf_path)
        pdf_name_without_ext = os.path.splitext(base_name)[0].replace(" ", "_")

        # Use the existing directory '../../ppt_references/'
        new_dir_path = os.path.join(os.getcwd(), "multimodal/ppt_references")

        image_paths = []

        for page_num, page in enumerate(doc.pages):
            pix = page.to_image()

            # Save images in the existing directory
            output_image_path = os.path.join(new_dir_path, f"{pdf_name_without_ext}_{page_num:04d}.png")
            pix.save(output_image_path)
            image_paths.append((output_image_path, page_num))

        doc.close()
        return image_paths
    
    def extract_text_and_notes_from_ppt(ppt_path):
        prs = Presentation(ppt_path)
        text_and_notes = []
        for slide in prs.slides:
            slide_text = ' '.join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
            try:
                notes = slide.notes_slide.notes_text_frame.text if slide.notes_slide else ""
            except:
                notes = ""
            text_and_notes.append((slide_text, notes))
        return text_and_notes